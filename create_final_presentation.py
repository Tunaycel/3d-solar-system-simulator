"""
Generate Final Project Presentation (PPTX)
==========================================
Run this script to produce: 3D_Graphics_Final_Presentation.pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUT_PATH = r"C:\Users\Tunay\.gemini\antigravity\scratch\3d_graphics_final_project\3D_Graphics_Final_Presentation.pptx"

# ── Colour palette ──────────────────────────────────────────────
BG_DARK   = RGBColor(0x05, 0x05, 0x18)   # deep navy
GOLD      = RGBColor(0xFF, 0xC0, 0x30)   # solar gold
CYAN      = RGBColor(0x30, 0xD0, 0xFF)   # neon cyan
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY= RGBColor(0xCC, 0xCC, 0xDD)
DARK_GREY = RGBColor(0x20, 0x20, 0x38)


def set_slide_bg(slide, r, g, b):
    """Fill slide background with a solid colour."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(r, g, b)


def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf   = txBox.text_frame
    tf.word_wrap = True
    p    = tf.paragraphs[0]
    p.alignment = align
    run  = p.add_run()
    run.text = text
    run.font.size    = Pt(font_size)
    run.font.bold    = bold
    run.font.italic  = italic
    run.font.color.rgb = color
    return txBox


def add_bullet_box(slide, lines, left, top, width, height, font_size=16):
    """Add a multi-line bullet list textbox."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.color.rgb = LIGHT_GREY
    return txBox


def add_rect(slide, left, top, width, height, fill_rgb, alpha=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.fill.background()
    return shape


# ════════════════════════════════════════════════════════════════════
# SLIDES
# ════════════════════════════════════════════════════════════════════

def slide1_title(prs):
    """Title slide."""
    layout = prs.slide_layouts[6]   # blank
    slide  = prs.slides.add_slide(layout)
    set_slide_bg(slide, 0x05, 0x05, 0x18)

    # Accent bar top
    add_rect(slide, 0, 0, 10, 0.08, GOLD)

    # Decorative circle (simulated with thin rectangle overlap)
    add_rect(slide, 7.0, 0.5, 2.8, 5.8, RGBColor(0x10, 0x10, 0x30))

    # Title
    add_textbox(slide, "3D Solar System Simulator",
                0.4, 0.8, 6.5, 1.2,
                font_size=36, bold=True, color=GOLD, align=PP_ALIGN.LEFT)

    add_textbox(slide, "Final Project — 3D Graphics",
                0.4, 2.0, 6.0, 0.6,
                font_size=20, color=CYAN, align=PP_ALIGN.LEFT)

    # Divider
    add_rect(slide, 0.4, 2.7, 5.5, 0.04, CYAN)

    meta = [
        "Student :  Hüseyin Tunay Çelik",
        "ID      :  97294",
        "Course  :  3D Graphics",
        "Lecturer:  Vishnu Suresh",
        "Univ.   :  WSB Merito Wrocław",
    ]
    add_bullet_box(slide, meta, 0.4, 2.9, 5.8, 2.2, font_size=15)

    # Accent bar bottom
    add_rect(slide, 0, 7.0, 10, 0.08, GOLD)


def slide2_overview(prs):
    layout = prs.slide_layouts[6]
    slide  = prs.slides.add_slide(layout)
    set_slide_bg(slide, 0x05, 0x05, 0x18)
    add_rect(slide, 0, 0, 10, 0.08, GOLD)
    add_rect(slide, 0, 0.08, 10, 0.8, RGBColor(0x0A, 0x0A, 0x28))

    add_textbox(slide, "Project Overview", 0.4, 0.15, 9, 0.65,
                font_size=28, bold=True, color=GOLD)

    desc = (
        "The Solar System Simulator is a real-time interactive 3D application "
        "that renders the 8 planets of our solar system orbiting around the Sun, "
        "complete with a Moon orbiting Earth, Saturn's rings, and a dense cosmic starfield.\n\n"
        "Built with Python 3, PyOpenGL, and GLFW, the application demonstrates "
        "the core concepts covered throughout the 3D Graphics course at WSB Merito Wrocław."
    )
    add_textbox(slide, desc, 0.4, 1.1, 9.2, 1.8, font_size=15, color=LIGHT_GREY)

    # 3 condition boxes
    boxes = [
        ("✅  Condition 1", "At least 3 visible 3D objects",
         "Sun + 8 Planets + Moon\n+ Saturn's rings = 11 objects", 0.3),
        ("✅  Condition 2", "Meaningful background / scene",
         "600-star starfield, orbital\npaths, dynamic point lighting", 3.6),
        ("✅  Condition 3", "Some form of animation",
         "Orbital motion, axial\nrotation, interactive camera", 6.9),
    ]
    for title, sub, detail, lft in boxes:
        add_rect(slide, lft, 3.1, 3.0, 3.5, RGBColor(0x10, 0x10, 0x35))
        add_textbox(slide, title, lft+0.1, 3.2, 2.8, 0.5,
                    font_size=14, bold=True, color=CYAN)
        add_textbox(slide, sub, lft+0.1, 3.7, 2.8, 0.4,
                    font_size=12, bold=True, color=GOLD)
        add_textbox(slide, detail, lft+0.1, 4.1, 2.8, 1.0,
                    font_size=12, color=LIGHT_GREY)

    add_rect(slide, 0, 7.0, 10, 0.08, GOLD)


def slide3_tech(prs):
    layout = prs.slide_layouts[6]
    slide  = prs.slides.add_slide(layout)
    set_slide_bg(slide, 0x05, 0x05, 0x18)
    add_rect(slide, 0, 0, 10, 0.08, GOLD)
    add_rect(slide, 0, 0.08, 10, 0.8, RGBColor(0x0A, 0x0A, 0x28))

    add_textbox(slide, "Technical Foundation", 0.4, 0.15, 9, 0.65,
                font_size=28, bold=True, color=GOLD)

    left_items = [
        "🔵  Technology Stack",
        "   • Language  : Python 3.11",
        "   • Windowing : GLFW",
        "   • Renderer  : PyOpenGL (OpenGL 1.x)",
        "   • Math      : NumPy",
        "   • Quadrics  : GLU Sphere / Quadric",
        "",
        "🔵  Geometry",
        "   • All 3D objects drawn with gluSphere()",
        "   • Saturn rings: GL_TRIANGLE_STRIP",
        "   • Orbit paths: GL_LINE_LOOP",
        "   • Starfield : GL_POINTS (600 stars)",
    ]
    right_items = [
        "🔵  Lighting Model",
        "   • GL_LIGHT0 at Sun origin (point light)",
        "   • Diffuse + Specular per planet",
        "   • Attenuation: quadratic falloff",
        "   • GL_NORMALIZE for correct normals",
        "",
        "🔵  Camera System",
        "   • Spherical coordinates (yaw, pitch, dist)",
        "   • Mouse drag → rotate",
        "   • Scroll wheel → zoom",
        "   • gluLookAt() updates each frame",
    ]
    add_bullet_box(slide, left_items,  0.4, 1.1, 4.5, 5.5, font_size=13)
    add_bullet_box(slide, right_items, 5.1, 1.1, 4.5, 5.5, font_size=13)

    add_rect(slide, 0, 7.0, 10, 0.08, GOLD)


def slide4_animation(prs):
    layout = prs.slide_layouts[6]
    slide  = prs.slides.add_slide(layout)
    set_slide_bg(slide, 0x05, 0x05, 0x18)
    add_rect(slide, 0, 0, 10, 0.08, GOLD)
    add_rect(slide, 0, 0.08, 10, 0.8, RGBColor(0x0A, 0x0A, 0x28))

    add_textbox(slide, "Animation & Scene Design", 0.4, 0.15, 9, 0.65,
                font_size=28, bold=True, color=GOLD)

    rows = [
        ("Object",       "Orbit Speed", "Rotation",   "Special Feature"),
        ("☀  Sun",       "—",          "—",          "Glow layers (alpha blended)"),
        ("☿  Mercury",  "4.74°/s",    "10°/s",      "Smallest orbit"),
        ("♀  Venus",    "3.50°/s",    "6°/s",       "Warm tones"),
        ("🌍  Earth",   "2.98°/s",    "15°/s",      "Has orbiting Moon"),
        ("🌙  Moon",    "13.0°/s*",   "—",          "*around Earth"),
        ("♂  Mars",    "2.41°/s",    "14°/s",      "Red-orange tint"),
        ("♃  Jupiter", "1.31°/s",    "35°/s",      "Largest planet"),
        ("♄  Saturn",  "0.97°/s",    "32°/s",      "Tilted ring system"),
        ("⛢  Uranus",  "0.68°/s",    "20°/s",      "Ice-blue hue"),
        ("♆  Neptune", "0.54°/s",    "22°/s",      "Deepest blue"),
    ]

    col_x = [0.3, 3.3, 5.2, 7.0]
    for r_idx, row in enumerate(rows):
        y = 1.1 + r_idx * 0.54
        if r_idx == 0:
            clr = GOLD
            sz  = 13
            bld = True
        else:
            clr = LIGHT_GREY
            sz  = 12
            bld = False
        for c_idx, cell in enumerate(row):
            add_textbox(slide, cell, col_x[c_idx], y, 2.8, 0.5,
                        font_size=sz, bold=bld, color=clr)

    add_rect(slide, 0, 7.0, 10, 0.08, GOLD)


def slide5_code(prs):
    layout = prs.slide_layouts[6]
    slide  = prs.slides.add_slide(layout)
    set_slide_bg(slide, 0x05, 0x05, 0x18)
    add_rect(slide, 0, 0, 10, 0.08, GOLD)
    add_rect(slide, 0, 0.08, 10, 0.8, RGBColor(0x0A, 0x0A, 0x28))

    add_textbox(slide, "Key Code Concepts", 0.4, 0.15, 9, 0.65,
                font_size=28, bold=True, color=GOLD)

    snippets = [
        ("Lighting Setup (Sun as light source)",
         "glLightfv(GL_LIGHT0, GL_POSITION, [0,0,0,1.0])  # point at origin\n"
         "glLightfv(GL_LIGHT0, GL_DIFFUSE,  [1.0,1.0,0.9,1.0])\n"
         "glLightf (GL_LIGHT0, GL_QUADRATIC_ATTENUATION, 0.001)"),
        ("Orbital Motion (per-frame update)",
         "orbit_angle = (orbit_speed * sim_time) % 360.0\n"
         "ox = cos(radians(orbit_angle)) * orbit_radius\n"
         "oz = sin(radians(orbit_angle)) * orbit_radius\n"
         "glTranslatef(ox, 0.0, oz)"),
        ("Spherical Camera",
         "eye_x = dist * cos(pitch) * sin(yaw)\n"
         "eye_y = dist * sin(pitch)\n"
         "eye_z = dist * cos(pitch) * cos(yaw)\n"
         "gluLookAt(eye_x, eye_y, eye_z,  0,0,0,  0,1,0)"),
    ]

    for i, (title, code) in enumerate(snippets):
        y = 1.05 + i * 2.1
        add_rect(slide, 0.3, y, 9.4, 1.85, RGBColor(0x0C, 0x0C, 0x28))
        add_textbox(slide, f"▸  {title}", 0.4, y + 0.05, 9, 0.35,
                    font_size=13, bold=True, color=CYAN)
        add_textbox(slide, code, 0.5, y + 0.40, 9, 1.35,
                    font_size=11, color=GOLD)

    add_rect(slide, 0, 7.0, 10, 0.08, GOLD)


def slide6_instructions(prs):
    layout = prs.slide_layouts[6]
    slide  = prs.slides.add_slide(layout)
    set_slide_bg(slide, 0x05, 0x05, 0x18)
    add_rect(slide, 0, 0, 10, 0.08, GOLD)
    add_rect(slide, 0, 0.08, 10, 0.8, RGBColor(0x0A, 0x0A, 0x28))

    add_textbox(slide, "How to Run the Project", 0.4, 0.15, 9, 0.65,
                font_size=28, bold=True, color=GOLD)

    steps = [
        ("Step 1 — Install dependencies",
         "pip install pyopengl pyopengl-accelerate glfw numpy"),
        ("Step 2 — Run the application",
         "python main.py"),
        ("Step 3 — Interact with the simulation",
         "Mouse Left-Click + Drag  →  Rotate camera\n"
         "Scroll Wheel             →  Zoom in / out\n"
         "R key                    →  Reset camera to default view\n"
         "+ / - keys               →  Speed up / slow down simulation\n"
         "ESC                      →  Quit"),
    ]

    for i, (title, detail) in enumerate(steps):
        y = 1.1 + i * 1.9
        add_rect(slide, 0.3, y, 9.4, 1.7, RGBColor(0x0C, 0x0C, 0x28))
        add_textbox(slide, title, 0.45, y + 0.1, 9, 0.4,
                    font_size=15, bold=True, color=CYAN)
        add_textbox(slide, detail, 0.5, y + 0.5, 9, 1.1,
                    font_size=13, color=GOLD)

    add_textbox(slide, "Submitted by  Hüseyin Tunay Çelik  •  97294  •  WSB Merito Wrocław",
                0.3, 6.8, 9.4, 0.35, font_size=11, color=LIGHT_GREY, italic=True,
                align=PP_ALIGN.CENTER)
    add_rect(slide, 0, 7.0, 10, 0.08, GOLD)


# ════════════════════════════════════════════════════════════════════
def main():
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7.5)

    slide1_title     (prs)
    slide2_overview  (prs)
    slide3_tech      (prs)
    slide4_animation (prs)
    slide5_code      (prs)
    slide6_instructions(prs)

    prs.save(OUT_PATH)
    print(f"[OK] Presentation saved to: {OUT_PATH}")

if __name__ == "__main__":
    main()
